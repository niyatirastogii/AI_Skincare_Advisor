import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_cream = RGBColor(247, 245, 240)
    c_charcoal = RGBColor(30, 30, 30)
    c_gold = RGBColor(212, 175, 55)
    c_gray = RGBColor(100, 100, 100)
    c_white = RGBColor(255, 255, 255)
    
    # --- SLIDE 1: Title Slide (Dark Theme for Premium First Impression) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = c_charcoal
    bg1.line.fill.background()
    
    # Accent shape - thin gold rectangle on left
    accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = c_gold
    accent1.line.fill.background()
    
    # Title Text Box (Centered horizontally)
    tb_title = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.33), Inches(4.0))
    tf = tb_title.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "AI Skincare Advisor"
    p_title.font.name = 'Outfit'
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = c_gold
    p_title.alignment = PP_ALIGN.CENTER
    p_title.space_after = Pt(14)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Personalized Skincare Diagnostics and Recommendations Driven by Computer Vision and Generative AI"
    p_sub.font.name = 'Outfit'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = c_cream
    p_sub.font.italic = True
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_after = Pt(28)
    
    p_author = tf.add_paragraph()
    p_author.text = "Project Overview Presentation"
    p_author.font.name = 'Outfit'
    p_author.font.size = Pt(14)
    p_author.font.color.rgb = c_gray
    p_author.alignment = PP_ALIGN.CENTER
    
    # Helper for content slides (Light Theme)
    def add_content_slide(title, category="AI SKINCARE ADVISOR"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = c_cream
        bg.line.fill.background()
        
        # Gold header strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        strip.fill.solid()
        strip.fill.fore_color.rgb = c_gold
        strip.line.fill.background()
        
        # Category label
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = tb_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Outfit'
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_gold
        
        # Title
        tb_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.0))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Outfit'
        p_t.font.size = Pt(36)
        p_t.font.bold = True
        p_t.font.color.rgb = c_charcoal
        
        return slide

    # --- SLIDE 2: Project Overview & Objectives ---
    slide2 = add_content_slide("Project Overview & Objectives")
    
    # Left highlighted card
    left_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.5))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = c_charcoal
    left_card.line.fill.background()
    
    tb_lc = slide2.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(3.9), Inches(3.9))
    tf_lc = tb_lc.text_frame
    tf_lc.word_wrap = True
    
    p_lc_head = tf_lc.paragraphs[0]
    p_lc_head.text = "THE MISSION"
    p_lc_head.font.name = 'Outfit'
    p_lc_head.font.size = Pt(14)
    p_lc_head.font.bold = True
    p_lc_head.font.color.rgb = c_gold
    p_lc_head.space_after = Pt(14)
    
    p_lc_body = tf_lc.add_paragraph()
    p_lc_body.text = "Democratizing expert dermatological advice through an accessible, privacy-focused, and intelligent digital assistant. By fusing local computer vision with state-of-the-art generative AI, we deliver actionable skincare insights instantly."
    p_lc_body.font.name = 'Outfit'
    p_lc_body.font.size = Pt(16)
    p_lc_body.font.color.rgb = c_cream
    p_lc_body.line_spacing = 1.3
    
    # Right column bullets
    tb_rc = slide2.shapes.add_textbox(Inches(5.8), Inches(1.9), Inches(6.7), Inches(4.8))
    tf_rc = tb_rc.text_frame
    tf_rc.word_wrap = True
    
    bullets2 = [
        ("The Skincare Challenge", "Consumers struggle with identifying skin types and conditions, often wasting resources on incompatible products containing harmful active ingredients."),
        ("Hybrid Diagnostic Approach", "Combines deterministic Computer Vision (OpenCV/MediaPipe) for structural face mapping with probabilistic Generative AI (Gemini) for complex visual and contextual skin analysis."),
        ("Data-Driven Personalization", "Cross-references identified skin concerns with a localized, structured product database to offer objective recommendations based on ingredient compatibility.")
    ]
    
    for i, (head, text) in enumerate(bullets2):
        p_h = tf_rc.paragraphs[0] if i == 0 else tf_rc.add_paragraph()
        p_h.text = f"•  {head}"
        p_h.font.name = 'Outfit'
        p_h.font.size = Pt(18)
        p_h.font.bold = True
        p_h.font.color.rgb = c_charcoal
        p_h.space_after = Pt(2)
        if i > 0: p_h.space_before = Pt(12)
        
        p_b = tf_rc.add_paragraph()
        p_b.text = text
        p_b.font.name = 'Outfit'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = c_gray
        p_b.space_after = Pt(4)
        
    # --- SLIDE 3: Core Features ---
    slide3 = add_content_slide("Core Features & User Experience")
    
    features = [
        ("📸 Skin Analysis", "Extracts real-time features using MediaPipe Face Mesh. Computes facial redness, estimates acne density, and generates an automated skin health score with localized heatmaps."),
        ("🧴 Recommendations", "Searches a custom product catalog of thousands of skincare products. Recommends specific active ingredients (e.g., Niacinamide, Salicylic Acid) based on user's unique profile."),
        ("🌿 Home Remedies", "Provides natural, safe, and accessible home remedy instructions generated dynamically by Gemini. Tailors application methods and benefits to user's severity levels.")
    ]
    
    for idx, (title, desc) in enumerate(features):
        left_pos = Inches(0.8 + idx * 4.0)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.0), Inches(3.7), Inches(4.5))
        card.fill.solid()
        # Make middle card highlight with charcoal theme
        if idx == 1:
            card.fill.fore_color.rgb = c_charcoal
            card.line.fill.background()
            text_color_title = c_gold
            text_color_desc = c_cream
        else:
            card.fill.fore_color.rgb = c_white
            card.line.color.rgb = c_gold
            card.line.width = Pt(1)
            text_color_title = c_charcoal
            text_color_desc = c_gray
            
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), Inches(3.3), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Outfit'
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = text_color_title
        p_t.space_after = Pt(12)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Outfit'
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = text_color_desc
        p_d.line_spacing = 1.3
        
    # --- SLIDE 4: System Architecture ---
    slide4 = add_content_slide("System Architecture & Data Flow")
    
    # Left column: Description
    tb_arch_l = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(4.8), Inches(4.5))
    tf_arch_l = tb_arch_l.text_frame
    tf_arch_l.word_wrap = True
    
    p_arch_t = tf_arch_l.paragraphs[0]
    p_arch_t.text = "Multi-Layered Architecture"
    p_arch_t.font.name = 'Outfit'
    p_arch_t.font.size = Pt(22)
    p_arch_t.font.bold = True
    p_arch_t.font.color.rgb = c_charcoal
    p_arch_t.space_after = Pt(14)
    
    p_arch_b = tf_arch_l.add_paragraph()
    p_arch_b.text = "The application splits workload between local computation and remote models to optimize response times and privacy. Local Computer Vision handles immediate geometry and metrics, while LLM services perform deep semantic diagnostics.\n\nCustom caching ensures that repeated lookups in the product database do not degrade responsiveness."
    p_arch_b.font.name = 'Outfit'
    p_arch_b.font.size = Pt(14)
    p_arch_b.font.color.rgb = c_gray
    p_arch_b.line_spacing = 1.3
    
    # Right column: The 3 architectural layers as blocks
    layers = [
        ("1. Presentation Layer (Streamlit Dashboard)", "Features a dark gold glassmorphic user interface with interactive camera inputs, image uploads, and dynamic visualization widgets."),
        ("2. Analysis Engine (CV & LLM Hybrid)", "Local OpenCV/MediaPipe calculates redness/acne density and face mesh connections. Remote Gemini 2.5 Flash analyzes complex image context."),
        ("3. Product Database Layer (Pandas Queries)", "Loads and indexes 40MB+ of curated products. Performs custom text-matching searches over ingredient lists and usage types with relevance thresholds.")
    ]
    
    for i, (name, details) in enumerate(layers):
        top_pos = Inches(2.0 + i * 1.5)
        block = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), top_pos, Inches(6.5), Inches(1.3))
        block.fill.solid()
        block.fill.fore_color.rgb = c_white
        block.line.color.rgb = c_gold
        block.line.width = Pt(1)
        
        tb_b = slide4.shapes.add_textbox(Inches(6.1), top_pos + Inches(0.05), Inches(6.3), Inches(1.2))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        
        p_bn = tf_b.paragraphs[0]
        p_bn.text = name
        p_bn.font.name = 'Outfit'
        p_bn.font.size = Pt(15)
        p_bn.font.bold = True
        p_bn.font.color.rgb = c_charcoal
        p_bn.space_after = Pt(2)
        
        p_bd = tf_b.add_paragraph()
        p_bd.text = details
        p_bd.font.name = 'Outfit'
        p_bd.font.size = Pt(11)
        p_bd.font.color.rgb = c_gray
        
    # --- SLIDE 5: Technology Stack ---
    slide5 = add_content_slide("Technology Stack & APIs")
    
    # 4 grid blocks for technologies
    techs = [
        ("Streamlit", "Web Interface & State", "Powers the user interface, routing, styling, state management, and file upload systems with zero-latency re-renders."),
        ("OpenCV & MediaPipe", "Computer Vision Analysis", "Performs real-time face detection, mesh landmark extraction, redness scoring, and color-mapped visual heatmaps."),
        ("Gemini 2.5 Flash", "Generative AI Image Analysis", "Conducts semantic visual analysis, detects skin type, severity of concerns, recommends ingredients, and provides dermatologist advice."),
        ("Pandas & Custom Regex", "Product Database & Search", "Filters, clean-matches, and queries thousands of skincare records in a 40MB database by ingredient keywords and skin types.")
    ]
    
    for idx, (name, sub, desc) in enumerate(techs):
        row = idx // 2
        col = idx % 2
        
        left_pos = Inches(0.8 + col * 6.1)
        top_pos = Inches(2.0 + row * 2.3)
        
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.6), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = c_white
        card.line.color.rgb = c_gold
        card.line.width = Pt(1)
        
        tb = slide5.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.1), Inches(5.2), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_name = tf.paragraphs[0]
        p_name.text = name
        p_name.font.name = 'Outfit'
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = c_charcoal
        p_name.space_after = Pt(2)
        
        p_sub = tf.add_paragraph()
        p_sub.text = sub.upper()
        p_sub.font.name = 'Outfit'
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = c_gold
        p_sub.space_after = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Outfit'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = c_gray
        p_desc.line_spacing = 1.2
        
    # --- SLIDE 6: Summary & Future Scope (Dark Theme) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = c_charcoal
    bg6.line.fill.background()
    
    # Accent shape - thin gold rectangle on left
    accent6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    accent6.fill.solid()
    accent6.fill.fore_color.rgb = c_gold
    accent6.line.fill.background()
    
    tb_title6 = slide6.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.8))
    tf_t6 = tb_title6.text_frame
    p_t6 = tf_t6.paragraphs[0]
    p_t6.text = "Summary & Future Roadmap"
    p_t6.font.name = 'Outfit'
    p_t6.font.size = Pt(36)
    p_t6.font.bold = True
    p_t6.font.color.rgb = c_gold
    
    tb_c6 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
    tf_c6 = tb_c6.text_frame
    tf_c6.word_wrap = True
    
    points6 = [
        ("Project Summary", "AI Skincare Advisor successfully implements a dual computer-vision and LLM analysis dashboard to perform structural and semantic analysis, bridging the gap between raw data and dermatologist advice."),
        ("Routine Personalization & Tracking", "Future releases will integrate user-profile saving and routine tracking over time to visually display changes in redness and acne scores using graphs."),
        ("Clinical Validations & Expert Integrations", "Introduce confidence levels, direct links to skin doctors, and diagnostic reports to elevate the app into a fully validated clinical-support tool."),
        ("E-Commerce Integrations", "Partner with major skincare platforms to enable direct-to-cart purchasing of recommended products directly from the advisor dashboard.")
    ]
    
    for i, (head, text) in enumerate(points6):
        p_h = tf_c6.paragraphs[0] if i == 0 else tf_c6.add_paragraph()
        p_h.text = f"•  {head}"
        p_h.font.name = 'Outfit'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = c_gold
        p_h.space_after = Pt(2)
        if i > 0: p_h.space_before = Pt(10)
        
        p_b = tf_c6.add_paragraph()
        p_b.text = text
        p_b.font.name = 'Outfit'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_cream
        p_b.space_after = Pt(4)
        
    prs.save("AI_Skincare_Advisor_Presentation.pptx")
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_deck()
